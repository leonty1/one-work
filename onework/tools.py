"""onework 核心工具函数

基于以下开源项目构建：
- Microsoft markitdown: https://github.com/microsoft/markitdown
- mammoth: Word 文档解析
- pymupdf: PDF 处理
- python-pptx: PPT 处理
- trafilatura: 网页内容提取
- markdownify: HTML 转 Markdown

感谢所有开源贡献者！
"""

import os
import json
from typing import Optional, List, Dict, Any
from pathlib import Path

# 依赖可用性检查
_MARKITDOWN_AVAILABLE = None
_TRAFILATURA_AVAILABLE = None
_MARKDOWNIFY_AVAILABLE = None


def _check_markitdown() -> bool:
    """检查 markitdown 是否可用"""
    global _MARKITDOWN_AVAILABLE
    if _MARKITDOWN_AVAILABLE is None:
        try:
            import markitdown
            _MARKITDOWN_AVAILABLE = True
        except ImportError:
            _MARKITDOWN_AVAILABLE = False
    return _MARKITDOWN_AVAILABLE


def _check_trafilatura() -> bool:
    """检查 trafilatura 是否可用"""
    global _TRAFILATURA_AVAILABLE
    if _TRAFILATURA_AVAILABLE is None:
        try:
            import trafilatura
            _TRAFILATURA_AVAILABLE = True
        except ImportError:
            _TRAFILATURA_AVAILABLE = False
    return _TRAFILATURA_AVAILABLE


def _check_markdownify() -> bool:
    """检查 markdownify 是否可用"""
    global _MARKDOWNIFY_AVAILABLE
    if _MARKDOWNIFY_AVAILABLE is None:
        try:
            import markdownify
            _MARKDOWNIFY_AVAILABLE = True
        except ImportError:
            _MARKDOWNIFY_AVAILABLE = False
    return _MARKDOWNIFY_AVAILABLE


def read_document(file_path: str, use_markitdown: bool = True) -> Dict[str, Any]:
    """
    读取文档并返回结构化信息。
    
    自动识别格式（Word/PDF/PPT/图片），返回：
    - content: Markdown 格式全文
    - structure: 文档结构（标题层级、章节）
    - tables: 表格列表
    - file_type: 文件类型
    
    Args:
        file_path: 文档的绝对路径
        use_markitdown: 是否优先使用 markitdown (默认 True)
    
    Returns:
        包含 content, structure, tables, file_type 的字典
    """
    path = Path(file_path)
    
    if not path.exists():
        return {
            "success": False,
            "error": f"文件不存在: {file_path}"
        }
    
    ext = path.suffix.lower()
    
    # 优先使用 markitdown (如果可用且用户未禁用)
    if use_markitdown and _check_markitdown():
        md_result = _read_with_markitdown(file_path)
        if md_result.get("success"):
            return md_result
    
    # 回退到原有实现
    if ext == ".docx":
        return _read_docx(file_path)
    elif ext == ".pdf":
        return _read_pdf(file_path)
    elif ext == ".pptx":
        return _read_pptx(file_path)
    elif ext in [".md", ".markdown"]:
        return _read_markdown(file_path)
    elif ext in [".png", ".jpg", ".jpeg", ".gif", ".bmp"]:
        return _read_image(file_path)
    else:
        return {
            "success": False,
            "error": f"不支持的格式: {ext}"
        }


def _read_with_markitdown(file_path: str) -> Dict[str, Any]:
    """使用 markitdown 读取文档"""
    try:
        from markitdown import MarkItDown
        
        md = MarkItDown()
        result = md.convert(file_path)
        content = result.text_content
        
        structure = _extract_headings_from_md(content)
        
        return {
            "success": True,
            "content": content,
            "structure": structure,
            "tables": _extract_tables_from_md(content),
            "file_type": Path(file_path).suffix.lower().lstrip(".")
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"markitdown 转换失败: {str(e)}"
        }


def read_url(url: str, include_links: bool = True, include_images: bool = True) -> Dict[str, Any]:
    """
    读取网页 URL 并转换为 Markdown。
    
    自动提取网页正文内容，过滤广告和导航，返回：
    - content: Markdown 格式全文
    - title: 网页标题
    - author: 作者（如有）
    - date: 发布日期（如有）
    - url: 原始 URL
    
    Args:
        url: 网页 URL
        include_links: 是否保留链接（默认 True）
        include_images: 是否包含图片（默认 True）
    
    Returns:
        包含 content, title, author, date, url 的字典
    """
    # 优先使用 trafilatura（智能提取正文）
    if _check_trafilatura():
        return _read_url_with_trafilatura(url, include_links, include_images)
    
    # 回退到 markdownify + requests
    elif _check_markdownify():
        return _read_url_with_markdownify(url, include_links, include_images)
    
    else:
        return {
            "success": False,
            "error": "缺少网页解析依赖。请安装: pip install trafilatura markdownify requests"
        }


def _read_url_with_trafilatura(url: str, include_links: bool = True, include_images: bool = True) -> Dict[str, Any]:
    """使用 trafilatura 读取网页并转换为 Markdown"""
    try:
        import trafilatura
        from trafilatura.settings import use_config
        
        # 配置：包含链接和图片
        config = use_config()
        config.set("extractors", "links", str(include_links).lower())
        config.set("extractors", "images", str(include_images).lower())
        
        # 下载并提取内容
        downloaded = trafilatura.fetch_url(url)
        if not downloaded:
            return {
                "success": False,
                "error": f"无法下载网页: {url}"
            }
        
        # 提取元数据和内容
        result = trafilatura.bare_extraction(
            downloaded,
            config=config,
            include_links=include_links,
            include_images=include_images,
            output_format="json"
        )
        
        if not result:
            return {
                "success": False,
                "error": "无法提取网页内容"
            }
        
        # 转换为 Markdown
        content = trafilatura.extract(downloaded, config=config, include_links=include_links)
        
        # 解析标题结构
        structure = _extract_headings_from_md(content) if content else []
        
        return {
            "success": True,
            "content": content or "",
            "title": result.get("title", ""),
            "author": result.get("author", ""),
            "date": result.get("date", ""),
            "url": url,
            "structure": structure,
            "file_type": "url"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"trafilatura 转换失败: {str(e)}"
        }


def _read_url_with_markdownify(url: str, include_links: bool = True, include_images: bool = True) -> Dict[str, Any]:
    """使用 markdownify + requests 读取网页并转换为 Markdown"""
    try:
        import requests
        from markdownify import markdownify as md
        from bs4 import BeautifulSoup
        
        # 设置请求头
        headers = {
            "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
        }
        
        response = requests.get(url, headers=headers, timeout=30)
        response.raise_for_status()
        response.encoding = response.apparent_encoding or 'utf-8'
        
        html = response.text
        
        # 解析标题
        soup = BeautifulSoup(html, 'html.parser')
        title = ""
        if soup.title:
            title = soup.title.get_text(strip=True)
        elif soup.find('h1'):
            title = soup.find('h1').get_text(strip=True)
        
        # 尝试提取正文（简单启发式方法）
        main_content = None
        for selector in ['article', 'main', '.content', '#content', '.post', '.entry']:
            main_content = soup.select_one(selector)
            if main_content:
                break
        
        # 如果没有找到，使用 body
        if not main_content:
            main_content = soup.body or soup
        
        # 转换为 Markdown
        content = md(str(main_content), heading_style="ATX")
        
        # 解析标题结构
        structure = _extract_headings_from_md(content) if content else []
        
        return {
            "success": True,
            "content": content,
            "title": title,
            "author": "",
            "date": "",
            "url": url,
            "structure": structure,
            "file_type": "url"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"markdownify 转换失败: {str(e)}"
        }


def write_html(markdown_content: str, title: str = "Document", output_file: str = None) -> Dict[str, Any]:
    """
    将 Markdown 内容转换为 HTML 网页。
    
    Args:
        markdown_content: Markdown 格式内容
        title: 网页标题（默认 "Document"）
        output_file: 输出文件路径（可选，不指定则返回 HTML 字符串）
    
    Returns:
        success: 是否成功
        html: HTML 内容（当 output_file 未指定时返回）
        file: 输出文件路径（当 output_file 指定时返回）
    """
    try:
        import markdown
        
        # 配置 Markdown 扩展
        md = markdown.Markdown(
            extensions=[
                'extra',           # 表格、脚注等
                'codehilite',     # 代码高亮
                'toc',            # 目录
                'meta',           # 元数据
                'nl2br',          # 换行转 <br>
            ]
        )
        
        body_html = md.convert(markdown_content)
        
        # 构建完整 HTML
        html = f"""<!DOCTYPE html>
<html lang="zh-CN">
<head>
    <meta charset="UTF-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    <title>{title}</title>
    <style>
        body {{
            font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, "Helvetica Neue", Arial, sans-serif;
            line-height: 1.6;
            max-width: 900px;
            margin: 0 auto;
            padding: 20px;
            color: #333;
            background: #f5f5f5;
        }}
        .container {{
            background: white;
            padding: 40px;
            border-radius: 8px;
            box-shadow: 0 2px 10px rgba(0,0,0,0.1);
        }}
        h1, h2, h3, h4, h5, h6 {{
            color: #2c3e50;
            margin-top: 1.5em;
            margin-bottom: 0.5em;
        }}
        h1 {{ font-size: 2em; border-bottom: 2px solid #3498db; padding-bottom: 10px; }}
        h2 {{ font-size: 1.5em; border-bottom: 1px solid #ecf0f1; padding-bottom: 8px; }}
        code {{
            background: #f8f9fa;
            padding: 2px 6px;
            border-radius: 4px;
            font-family: "SFMono-Regular", Consolas, "Liberation Mono", Menlo, monospace;
            font-size: 0.9em;
        }}
        pre {{
            background: #f8f9fa;
            padding: 16px;
            border-radius: 6px;
            overflow-x: auto;
        }}
        pre code {{
            background: none;
            padding: 0;
        }}
        table {{
            border-collapse: collapse;
            width: 100%;
            margin: 1em 0;
        }}
        th, td {{
            border: 1px solid #ddd;
            padding: 12px;
            text-align: left;
        }}
        th {{
            background: #3498db;
            color: white;
        }}
        tr:nth-child(even) {{
            background: #f8f9fa;
        }}
        blockquote {{
            border-left: 4px solid #3498db;
            margin: 1em 0;
            padding-left: 16px;
            color: #7f8c8d;
        }}
        a {{
            color: #3498db;
            text-decoration: none;
        }}
        a:hover {{
            text-decoration: underline;
        }}
        img {{
            max-width: 100%;
            height: auto;
            border-radius: 4px;
        }}
        hr {{
            border: none;
            border-top: 1px solid #ecf0f1;
            margin: 2em 0;
        }}
    </style>
</head>
<body>
    <div class="container">
        {body_html}
    </div>
</body>
</html>"""
        
        if output_file:
            with open(output_file, 'w', encoding='utf-8') as f:
                f.write(html)
            return {
                "success": True,
                "file": output_file
            }
        else:
            return {
                "success": True,
                "html": html
            }
    except Exception as e:
        return {
            "success": False,
            "error": f"Markdown 转 HTML 失败: {str(e)}"
        }


def _read_docx(file_path: str) -> Dict[str, Any]:
    """读取 Word 文档"""
    try:
        import mammoth
        
        with open(file_path, "rb") as doc_file:
            result = mammoth.convert_to_markdown(doc_file)
            content = result.value
        
        structure = _extract_headings_from_md(content)
        
        return {
            "success": True,
            "content": content,
            "structure": structure,
            "tables": _extract_tables_from_md(content),
            "file_type": "docx"
        }
    except ImportError:
        return {
            "success": False,
            "error": "请安装 mammoth: pip install mammoth"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"读取 Word 失败: {str(e)}"
        }


def _read_pdf(file_path: str) -> Dict[str, Any]:
    """读取 PDF 文档"""
    try:
        import pymupdf
        
        doc = pymupdf.open(file_path)
        content_parts = []
        
        for page_num, page in enumerate(doc):
            text = page.get_text()
            content_parts.append(f"## 第 {page_num + 1} 页\n\n{text}")
        
        content = "\n\n---\n\n".join(content_parts)
        
        return {
            "success": True,
            "content": content,
            "structure": _extract_headings_from_md(content),
            "tables": [],
            "file_type": "pdf"
        }
    except ImportError:
        return {
            "success": False,
            "error": "请安装 pymupdf: pip install pymupdf"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"读取 PDF 失败: {str(e)}"
        }


def _read_pptx(file_path: str) -> Dict[str, Any]:
    """读取 PPT 文档"""
    try:
        from pptx import Presentation
        
        prs = Presentation(file_path)
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)
            
            if slide_text:
                slides_content.append(f"## 幻灯片 {i + 1}\n\n" + "\n".join(slide_text))
        
        content = "\n\n---\n\n".join(slides_content)
        
        return {
            "success": True,
            "content": content,
            "structure": [{"level": 1, "title": f"幻灯片 {i+1}", "page": i+1} for i in range(len(slides_content))],
            "tables": [],
            "file_type": "pptx"
        }
    except ImportError:
        return {
            "success": False,
            "error": "请安装 python-pptx: pip install python-pptx"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"读取 PPT 失败: {str(e)}"
        }


def _read_markdown(file_path: str) -> Dict[str, Any]:
    """读取 Markdown 文档"""
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        
        structure = _extract_headings_from_md(content)
        
        return {
            "success": True,
            "content": content,
            "structure": structure,
            "tables": _extract_tables_from_md(content),
            "file_type": "markdown"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"读取 Markdown 失败: {str(e)}"
        }


def _read_image(file_path: str) -> Dict[str, Any]:
    """读取图片（OCR）"""
    try:
        from PIL import Image
        import pytesseract
        
        image = Image.open(file_path)
        content = pytesseract.image_to_string(image, lang='chi_sim+eng')
        
        return {
            "success": True,
            "content": content,
            "structure": [],
            "tables": [],
            "file_type": "image"
        }
    except ImportError:
        return {
            "success": False,
            "error": "请安装 Pillow 和 pytesseract: pip install pillow pytesseract"
        }
    except Exception as e:
        return {
            "success": False,
            "error": f"OCR 识别失败: {str(e)}"
        }


def _extract_headings_from_md(content: str) -> List[Dict[str, Any]]:
    """从 Markdown 中提取标题结构"""
    import re
    
    headings = []
    for line in content.split('\n'):
        match = re.match(r'^(#{1,6})\s+(.+)$', line)
        if match:
            level = len(match.group(1))
            title = match.group(2).strip()
            headings.append({"level": level, "title": title})
    
    return headings


def _extract_tables_from_md(content: str) -> List[str]:
    """从 Markdown 中提取表格"""
    import re
    
    tables = []
    in_table = False
    current_table = []
    
    for line in content.split('\n'):
        if '|' in line:
            in_table = True
            current_table.append(line)
        elif in_table and current_table:
            tables.append('\n'.join(current_table))
            current_table = []
            in_table = False
    
    if current_table:
        tables.append('\n'.join(current_table))
    
    return tables


def extract_structure(file_path: str) -> List[Dict[str, Any]]:
    """
    提取文档的大纲结构（章节层级）。
    
    Args:
        file_path: 文档路径
    
    Returns:
        层级化的章节列表
    """
    result = read_document(file_path)
    
    if result.get("success"):
        return result.get("structure", [])
    else:
        return []


def update_section(
    file_path: str,
    section_path: str,
    new_content: str,
    output_path: Optional[str] = None
) -> Dict[str, Any]:
    """
    更新文档中的特定章节。
    
    Args:
        file_path: 原文档路径
        section_path: 章节路径，如 "1. 项目背景/1.1 概述"
        new_content: 新的 Markdown 内容
        output_path: 输出路径（可选，默认生成新文件）
    
    Returns:
        包含 success, output_path, summary 的字典
    """
    if not output_path:
        path = Path(file_path)
        output_path = str(path.parent / f"{path.stem}_v2{path.suffix}")
    
    result = read_document(file_path)
    
    if not result.get("success"):
        return result
    
    # 简单实现：生成新的 Markdown 文件（带章节标记）
    # 后续可扩展为直接生成 docx/pptx
    output_md_path = output_path.replace(path.suffix, ".md")
    
    with open(output_md_path, "w", encoding="utf-8") as f:
        f.write(f"# 文档更新\n\n")
        f.write(f"**更新章节**: {section_path}\n\n")
        f.write(f"---\n\n")
        f.write(new_content)
    
    return {
        "success": True,
        "output_path": output_md_path,
        "summary": f"已更新章节 '{section_path}'，生成新文档: {output_md_path}"
    }


def convert_format(
    input_path: str,
    output_format: str,
    output_path: Optional[str] = None,
    template: Optional[str] = None
) -> Dict[str, Any]:
    """
    转换文档格式。
    
    支持格式：docx, pdf, pptx, html, md
    
    Args:
        input_path: 输入文件路径
        output_format: 目标格式
        output_path: 输出路径（可选）
        template: 模板名称（可选）
    
    Returns:
        包含 success, output_path 的字典
    """
    input_path_obj = Path(input_path)
    
    if not input_path_obj.exists():
        return {
            "success": False,
            "error": f"输入文件不存在: {input_path}"
        }
    
    if not output_path:
        output_path = str(input_path_obj.parent / f"{input_path_obj.stem}_converted.{output_format}")
    
    # 先转换为 Markdown 中间格式
    read_result = read_document(input_path)
    
    if not read_result.get("success"):
        return read_result
    
    markdown_content = read_result.get("content", "")
    
    # 根据目标格式渲染
    if output_format == "md":
        # 直接输出 Markdown
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(markdown_content)
        
    elif output_format == "html":
        # Markdown -> HTML
        try:
            import markdown
            
            html = markdown.markdown(
                markdown_content,
                extensions=['tables', 'fenced_code', 'toc']
            )
            
            full_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <title>Document</title>
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; }}
        table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
        th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
        th {{ background: #f5f5f5; }}
        h1, h2, h3 {{ color: #333; }}
        pre {{ background: #f5f5f5; padding: 10px; overflow-x: auto; }}
    </style>
</head>
<body>
{html}
</body>
</html>"""
            
            with open(output_path, "w", encoding="utf-8") as f:
                f.write(full_html)
        except ImportError:
            return {"success": False, "error": "请安装 markdown: pip install markdown"}
    
    elif output_format == "docx":
        # Markdown -> Word
        try:
            from docx import Document
            from docx.shared import Pt, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            doc = Document()
            
            # 简单的 Markdown -> Word 转换
            lines = markdown_content.split('\n')
            in_code_block = False
            
            for line in lines:
                if line.startswith('```'):
                    in_code_block = not in_code_block
                    continue
                
                if in_code_block:
                    doc.add_paragraph(line, style='Code')
                    continue
                
                if line.startswith('#'):
                    # 标题
                    level = len(line) - len(line.lstrip('#'))
                    heading = doc.add_heading(line.lstrip('# '), level=level)
                    heading.alignment = WD_ALIGN_PARAGRAPH.LEFT
                elif line.strip():
                    # 段落
                    doc.add_paragraph(line)
                else:
                    # 空行
                    doc.add_paragraph()
            
            doc.save(output_path)
        except ImportError:
            return {"success": False, "error": "请安装 python-docx: pip install python-docx"}
    
    elif output_format == "pdf":
        # Markdown -> PDF (通过 HTML)
        html_path = output_path.replace('.pdf', '.html')
        
        # 先转换为 HTML
        try:
            import markdown
            html = markdown.markdown(
                markdown_content,
                extensions=['tables', 'fenced_code']
            )
            
            full_html = f"""<!DOCTYPE html>
<html>
<head>
    <meta charset="utf-8">
    <style>
        body {{ font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', sans-serif; max-width: 800px; margin: 0 auto; padding: 40px; }}
        table {{ border-collapse: collapse; width: 100%; margin: 1em 0; }}
        th, td {{ border: 1px solid #ddd; padding: 8px; text-align: left; }}
        h1, h2, h3 {{ color: #333; }}
    </style>
</head>
<body>{html}</body>
</html>"""
            
            with open(html_path, "w", encoding="utf-8") as f:
                f.write(full_html)
            
            # 使用 weasyprint 转换为 PDF
            try:
                from weasyprint import HTML
                HTML(html_path).write_pdf(output_path)
            except ImportError:
                return {
                    "success": True,
                    "output_path": html_path,
                    "message": "已转换为 HTML，PDF 需要安装 weasyprint: pip install weasyprint"
                }
        except ImportError:
            return {"success": False, "error": "请安装 markdown: pip install markdown"}
    
    elif output_format == "pptx":
        # Markdown -> PPT
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            prs = Presentation()
            
            # 简单的 Markdown -> PPT 转换
            # 按 H1 分页，每页一个标题 + 内容
            lines = markdown_content.split('\n')
            current_title = "封面"
            current_content = []
            
            for line in lines:
                if line.startswith('## '):
                    # 新的幻灯片
                    if current_content:
                        _add_slide(prs, current_title, current_content)
                    
                    current_title = line.replace('## ', '').strip()
                    current_content = []
                elif line.startswith('# '):
                    current_title = line.replace('# ', '').strip()
                elif line.strip():
                    current_content.append(line.strip())
            
            # 最后一页
            if current_content:
                _add_slide(prs, current_title, current_content)
            
            prs.save(output_path)
        except ImportError:
            return {"success": False, "error": "请安装 python-pptx: pip install python-pptx"}
    
    else:
        return {
            "success": False,
            "error": f"不支持的目标格式: {output_format}"
        }
    
    return {
        "success": True,
        "output_path": output_path
    }


def _add_slide(prs: Any, title: str, content: List[str]):
    """添加一张幻灯片"""
    slide_layout = prs.slide_layouts[1]  # 标题和内容
    slide = prs.slides.add_slide(slide_layout)
    
    title_elem = slide.shapes.title
    title_elem.text = title
    
    content_elem = slide.placeholders[1]
    text_frame = content_elem.text_frame
    text_frame.clear()
    
    for item in content[:10]:  # 最多 10 条
        p = text_frame.add_paragraph()
        p.text = item
        p.level = 0


def create_from_outline(
    outline: str,
    output_format: str,
    output_path: Optional[str] = None,
    template: Optional[str] = None
) -> Dict[str, Any]:
    """
    根据大纲生成完整文档。
    
    Args:
        outline: 文档大纲（Markdown 格式）
        output_format: 输出格式 (docx, pptx, pdf)
        output_path: 输出路径（可选）
        template: 模板名称（可选）
    
    Returns:
        包含 success, output_path 的字典
    """
    if not output_path:
        import tempfile
        output_path = os.path.join(tempfile.gettempdir(), f"onework_output.{output_format}")
    
    # 大纲本身就是 Markdown 内容
    # 可以直接用 convert_format 处理
    if output_format == "md":
        with open(output_path, "w", encoding="utf-8") as f:
            f.write(outline)
        
        return {
            "success": True,
            "output_path": output_path
        }
    
    # 其他格式需要先转为 Markdown 文件
    import tempfile
    with tempfile.NamedTemporaryFile(mode='w', suffix='.md', delete=False, encoding='utf-8') as f:
        f.write(outline)
        md_path = f.name
    
    # 调用 convert_format
    result = convert_format(md_path, output_format, output_path, template)
    
    # 清理临时文件
    try:
        os.unlink(md_path)
    except:
        pass
    
    return result


# ============ 模板相关函数 ============

def list_templates(format: str = None) -> List[str]:
    """
    列出可用的模板。
    
    Args:
        format: 格式过滤 (docx, pptx, pdf)
    
    Returns:
        模板文件列表
    """
    # 获取当前文件所在目录
    current_dir = Path(__file__).parent.parent
    
    templates_dir = current_dir / "templates"
    
    if not templates_dir.exists():
        return []
    
    templates = []
    
    if format:
        format_dir = templates_dir / format
        if format_dir.exists():
            templates = [str(f) for f in format_dir.glob("*") if f.is_file()]
    else:
        for subdir in templates_dir.iterdir():
            if subdir.is_dir():
                for f in subdir.glob("*"):
                    if f.is_file():
                        templates.append(str(f))
    
    return templates